import os
import uuid
import re
import shutil
import time
import urllib.parse
from fastapi import FastAPI, File, UploadFile, BackgroundTasks
from fastapi.responses import JSONResponse, FileResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from markitdown import MarkItDown
from pptx import Presentation
from openpyxl import load_workbook
from oletools.olevba import VBA_Parser
app = FastAPI(title="MarkItDown Web UI")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Ensure directories exist
os.makedirs("static/conversions", exist_ok=True)
app.mount("/static", StaticFiles(directory="static"), name="static")

md = MarkItDown()

def _extract_pivot_tables(wb_formula, wb_data):
    """Extract pivot tables by reading displayed cell values from pivot location range."""
    from openpyxl.utils import range_boundaries
    pivots_by_sheet = {}

    for ws_f in wb_formula.worksheets:
        if not ws_f._pivots:
            continue

        ws_d = wb_data[ws_f.title]
        sheet_pivots = []

        for pivot in ws_f._pivots:
            info = {"name": pivot.name or "Pivot Table"}

            # Source range metadata
            cache = pivot.cache
            if cache and cache.cacheSource and cache.cacheSource.worksheetSource:
                src = cache.cacheSource.worksheetSource
                info["source"] = f"{src.sheet}!{src.ref}" if src.sheet else str(src.ref or "")

            # Field names from cache
            if cache and cache.cacheFields:
                info["fields"] = [cf.name for cf in cache.cacheFields if cf.name]

            # Read cell values from pivot's rendered location
            if pivot.location and pivot.location.ref:
                ref = pivot.location.ref
                min_col, min_row, max_col, max_row = range_boundaries(ref)
                rows = []
                for r in range(min_row, max_row + 1):
                    row = []
                    for c in range(min_col, max_col + 1):
                        val = ws_d.cell(r, c).value
                        row.append(str(val) if val is not None else '')
                    rows.append(row)
                info["rows"] = rows

            sheet_pivots.append(info)

        if sheet_pivots:
            pivots_by_sheet[ws_f.title] = sheet_pivots

    return pivots_by_sheet


def _extract_vba_macros(file_path):
    """Extract VBA macros from Office files."""
    try:
        vba_parser = VBA_Parser(file_path)
        if not vba_parser.detect_vba_macros():
            vba_parser.close()
            return ""

        parts = ["\n## VBA Macros\n"]
        for (filename, stream_path, vba_filename, vba_code_chunk) in vba_parser.extract_macros():
            parts.append(f"### Module: {vba_filename}\n```vba\n{vba_code_chunk}\n```\n")
        
        vba_parser.close()
        return "\n".join(parts)
    except Exception as e:
        print(f"Error extracting VBA: {e}")
        return ""

def _escape_cell(text):
    """Escape pipe and newline for Markdown table cell."""
    return text.replace('|', '\\|').replace('\n', ' ')


def convert_excel_with_formulas(file_path):
    """Convert Excel (.xlsx) showing values, formulas, and pivot tables."""
    # Use read_only=True to prevent OOM on large files (streams XML instead of loading all objects into memory)
    wb_data = load_workbook(file_path, data_only=True, read_only=True)
    wb_formula = load_workbook(file_path, data_only=False, read_only=True)
    parts = []

    # Extract pivot tables (might be skipped if read_only=True doesn't support _pivots, but we gracefully catch it)
    try:
        pivots_by_sheet = _extract_pivot_tables(wb_formula, wb_data)
    except Exception:
        pivots_by_sheet = {}

    for name in wb_formula.sheetnames:
        ws_d, ws_f = wb_data[name], wb_formula[name]

        parts.append(f"## {name}\n")
        rows = []
        
        # Use iter_rows to stream row by row simultaneously (O(N) time and O(1) memory)
        for row_d, row_f in zip(ws_d.iter_rows(values_only=True), ws_f.iter_rows(values_only=False)):
            row_data = []
            for val, cell_f in zip(row_d, row_f):
                raw = cell_f.value
                if isinstance(raw, str) and raw.startswith('='):
                    text = f"{val if val is not None else ''} (`{raw}`)"
                else:
                    text = str(val) if val is not None else ''
                row_data.append(_escape_cell(text))
            
            # Only append if row has at least one non-empty cell (to avoid massive trailing blank rows)
            if any(c.strip() != '' for c in row_data):
                rows.append(row_data)

        if not rows:
            continue

        ncols = max(len(r) for r in rows) if rows else 0
        if ncols == 0:
            continue
            
        parts.append('| ' + ' | '.join((rows[0] + [''] * ncols)[:ncols]) + ' |')
        parts.append('| ' + ' | '.join(['---'] * ncols) + ' |')
        for row in rows[1:]:
            parts.append('| ' + ' | '.join((row + [''] * ncols)[:ncols]) + ' |')
        parts.append('')

        # Render pivot tables for this sheet
        if name in pivots_by_sheet:
            for pv in pivots_by_sheet[name]:
                parts.append(f"### Pivot Table: {pv['name']}\n")
                if pv.get("source"):
                    parts.append(f"**Source:** `{pv['source']}`\n")
                if pv.get("fields"):
                    parts.append(f"**Fields:** {', '.join(pv['fields'])}\n")

                pv_rows = pv.get("rows", [])
                if not pv_rows:
                    parts.append("*Pivot table detected but no cached data available.*\n")
                    continue

                ncols_pv = max(len(r) for r in pv_rows)
                # First row as header
                h_row = [_escape_cell(c) for c in (pv_rows[0] + [''] * ncols_pv)[:ncols_pv]]
                parts.append('| ' + ' | '.join(h_row) + ' |')
                parts.append('| ' + ' | '.join(['---'] * ncols_pv) + ' |')
                for row in pv_rows[1:]:
                    cells = [_escape_cell(c) for c in (row + [''] * ncols_pv)[:ncols_pv]]
                    parts.append('| ' + ' | '.join(cells) + ' |')
                parts.append('')

    wb_data.close()
    wb_formula.close()
    return '\n'.join(parts)

def cleanup_old_jobs(directory="static/conversions", max_size_mb=100):
    """Giữ thư mục conversions luôn dưới mức max_size_mb, xóa các job cũ nhất nếu vượt quá."""
    try:
        job_dirs = []
        total_size = 0
        for entry in os.scandir(directory):
            if entry.is_dir():
                dir_size = sum(f.stat().st_size for f in os.scandir(entry.path) if f.is_file())
                job_dirs.append({"path": entry.path, "mtime": entry.stat().st_mtime, "size": dir_size})
                total_size += dir_size
                
        if total_size > max_size_mb * 1024 * 1024:
            job_dirs.sort(key=lambda x: x["mtime"]) # Cũ nhất lên đầu
            for job in job_dirs:
                shutil.rmtree(job["path"], ignore_errors=True)
                total_size -= job["size"]
                # Cắt xuống còn 80% giới hạn để có khoảng trống
                if total_size <= (max_size_mb * 0.8) * 1024 * 1024:
                    break
    except Exception as e:
        print(f"Lỗi khi dọn dẹp thư mục: {e}")

@app.head("/")
@app.get("/")
def serve_index():
    return FileResponse("static/index.html")

def process_single_file(job_id: str, job_dir: str, file_path: str, filename: str):
    try:
        # Excel: extract values + formulas via openpyxl
        if filename.lower().endswith(('.xlsx', '.xlsm', '.xltx', '.xltm')):
            markdown_text = convert_excel_with_formulas(file_path)
        else:
            result = md.convert(file_path)
            markdown_text = result.text_content
            markdown_text = re.sub(r'\bNaN\b', '', markdown_text)
            markdown_text = re.sub(r'Unnamed:\s*\d+', '', markdown_text)

        # Extract VBA macros for Office files
        office_exts = ('doc', 'docx', 'docm', 'dot', 'dotx', 'dotm', 
                       'xls', 'xlsx', 'xlsm', 'xlsb', 'xlt', 'xltx', 'xltm', 
                       'ppt', 'pptx', 'pptm', 'pot', 'potx', 'potm', 'pps', 'ppsx', 'ppsm')
        if filename.lower().endswith(office_exts):
            vba_text = _extract_vba_macros(file_path)
            if vba_text:
                markdown_text += vba_text
        
        # PPTX specific logic to extract images and rewrite markdown
        if filename.lower().endswith('.pptx'):
            try:
                prs = Presentation(file_path)
                img_counter = {}
                for slide in prs.slides:
                    sorted_shapes = sorted(
                        slide.shapes,
                        key=lambda s: (
                            float("-inf") if not s.top else s.top,
                            float("-inf") if not s.left else s.left,
                        ),
                    )
                    for shape in sorted_shapes:
                        if not (shape.shape_type == 13 or (shape.shape_type == 14 and hasattr(shape, "image"))):
                            continue
                        try:
                            blob = shape.image.blob
                            ext = shape.image.ext
                        except Exception:
                            continue

                        placeholder_name = re.sub(r"\W", "", shape.name)
                        placeholder = f"{placeholder_name}.jpg"

                        count = img_counter.get(placeholder_name, 0)
                        img_counter[placeholder_name] = count + 1
                        save_name = f"{placeholder_name}_{count}.{ext}" if count > 0 else f"{placeholder_name}.{ext}"

                        img_path = os.path.join(job_dir, save_name)
                        with open(img_path, "wb") as img_f:
                            img_f.write(blob)

                        encoded_filename = urllib.parse.quote(save_name)
                        img_url = f"/static/conversions/{job_id}/{encoded_filename}"

                        old_ref = f"]({placeholder})"
                        new_ref = f"]({img_url})"
                        markdown_text = markdown_text.replace(old_ref, new_ref, 1)

            except Exception as e:
                print(f"Error extracting PPTX images: {e}")

        # Save the final markdown to a file for download
        md_filename = f"{os.path.splitext(filename)[0]}.md"
        md_filepath = os.path.join(job_dir, md_filename)
        with open(md_filepath, "w", encoding="utf-8") as f:
            f.write(markdown_text)
            
        # Write success flag
        with open(os.path.join(job_dir, "success.txt"), "w") as f:
            f.write(md_filename)
            
    except Exception as e:
        with open(os.path.join(job_dir, "error.txt"), "w") as f:
            f.write(str(e))
    finally:
        if os.path.exists(file_path):
            os.remove(file_path)

@app.post("/api/convert")
def convert_file(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    # Auto clean up before processing new file
    cleanup_old_jobs()
    
    job_id = str(uuid.uuid4())
    job_dir = f"static/conversions/{job_id}"
    os.makedirs(job_dir, exist_ok=True)
    
    file_path = os.path.join(job_dir, file.filename)
    
    try:
        # Save uploaded file
        with open(file_path, "wb") as f:
            shutil.copyfileobj(file.file, f)
            
        background_tasks.add_task(process_single_file, job_id, job_dir, file_path, file.filename)
            
        return {
            "success": True,
            "status": "processing",
            "job_id": job_id
        }
        
    except Exception as e:
        return JSONResponse(status_code=500, content={"success": False, "error": str(e)})

from typing import List

import json

def process_batch_files_task(job_id: str, job_dir: str, file_paths: list, filenames: list):
    try:
        results_data = []
        for i, file_path in enumerate(file_paths):
            filename = filenames[i]
            try:
                if filename.lower().endswith(('.xlsx', '.xlsm', '.xltx', '.xltm')):
                    markdown_text = convert_excel_with_formulas(file_path)
                else:
                    result = md.convert(file_path)
                    markdown_text = result.text_content
                    markdown_text = re.sub(r'\bNaN\b', '', markdown_text)
                    markdown_text = re.sub(r'Unnamed:\s*\d+', '', markdown_text)
                
                # Extract VBA macros for Office files
                office_exts = ('doc', 'docx', 'docm', 'dot', 'dotx', 'dotm', 
                               'xls', 'xlsx', 'xlsm', 'xlsb', 'xlt', 'xltx', 'xltm', 
                               'ppt', 'pptx', 'pptm', 'pot', 'potx', 'potm', 'pps', 'ppsx', 'ppsm')
                if filename.lower().endswith(office_exts):
                    vba_text = _extract_vba_macros(file_path)
                    if vba_text:
                        markdown_text += vba_text
                
                if filename.lower().endswith('.pptx'):
                    try:
                        prs = Presentation(file_path)
                        file_prefix = os.path.splitext(filename)[0]
                        img_counter = {}
                        for slide in prs.slides:
                            sorted_shapes = sorted(
                                slide.shapes,
                                key=lambda s: (
                                    float("-inf") if not s.top else s.top,
                                    float("-inf") if not s.left else s.left,
                                ),
                            )
                            for shape in sorted_shapes:
                                if not (shape.shape_type == 13 or (shape.shape_type == 14 and hasattr(shape, "image"))):
                                    continue
                                try:
                                    blob = shape.image.blob
                                    ext = shape.image.ext
                                except Exception:
                                    continue

                                placeholder_name = re.sub(r"\W", "", shape.name)
                                placeholder = f"{placeholder_name}.jpg"

                                count = img_counter.get(placeholder_name, 0)
                                img_counter[placeholder_name] = count + 1
                                save_name = f"{file_prefix}_{placeholder_name}_{count}.{ext}" if count > 0 else f"{file_prefix}_{placeholder_name}.{ext}"

                                img_path = os.path.join(job_dir, save_name)
                                with open(img_path, "wb") as img_f:
                                    img_f.write(blob)

                                encoded_filename = urllib.parse.quote(save_name)
                                img_url = f"/static/conversions/{job_id}/{encoded_filename}"

                                old_ref = f"]({placeholder})"
                                new_ref = f"]({img_url})"
                                markdown_text = markdown_text.replace(old_ref, new_ref, 1)

                    except Exception as e:
                        print(f"Error extracting PPTX images in batch: {e}")

                # Save md file
                md_filename = f"{os.path.splitext(filename)[0]}.md"
                with open(os.path.join(job_dir, md_filename), "w", encoding="utf-8") as f:
                    f.write(markdown_text)
                    
                results_data.append({
                    "filename": filename,
                    "markdown": markdown_text
                })
            except Exception as e:
                print(f"Error converting {filename}: {e}")
            finally:
                if os.path.exists(file_path):
                    os.remove(file_path)
            
        # Create ZIP archive
        zip_path = f"static/conversions/{job_id}_archive"
        shutil.make_archive(zip_path, 'zip', job_dir)
        
        # Write success data
        with open(os.path.join(job_dir, "success.json"), "w", encoding="utf-8") as f:
            json.dump(results_data, f)
            
    except Exception as e:
        with open(os.path.join(job_dir, "error.txt"), "w") as f:
            f.write(str(e))

@app.post("/api/convert_batch")
def convert_batch(background_tasks: BackgroundTasks, files: List[UploadFile] = File(...)):
    cleanup_old_jobs()
    
    job_id = str(uuid.uuid4())
    job_dir = f"static/conversions/{job_id}"
    os.makedirs(job_dir, exist_ok=True)
    
    try:
        file_paths = []
        filenames = []
        for file in files:
            file_path = os.path.join(job_dir, file.filename)
            with open(file_path, "wb") as f:
                shutil.copyfileobj(file.file, f)
            file_paths.append(file_path)
            filenames.append(file.filename)
            
        background_tasks.add_task(process_batch_files_task, job_id, job_dir, file_paths, filenames)
        
        return {
            "success": True,
            "status": "processing",
            "job_id": job_id
        }
        
    except Exception as e:
        return JSONResponse(status_code=500, content={"success": False, "error": str(e)})

@app.get("/api/status/{job_id}")
def check_status(job_id: str):
    job_dir = f"static/conversions/{job_id}"
    if not os.path.exists(job_dir):
        return {"success": False, "error": "Job not found"}
        
    if os.path.exists(os.path.join(job_dir, "error.txt")):
        with open(os.path.join(job_dir, "error.txt"), "r") as f:
            return {"success": False, "error": f.read()}
            
    # Check if single file success
    if os.path.exists(os.path.join(job_dir, "success.txt")):
        with open(os.path.join(job_dir, "success.txt"), "r") as f:
            md_filename = f.read().strip()
        md_filepath = os.path.join(job_dir, md_filename)
        with open(md_filepath, "r", encoding="utf-8") as f:
            markdown_text = f.read()
        return {
            "success": True,
            "status": "completed",
            "markdown": markdown_text,
            "download_url": f"/static/conversions/{job_id}/{md_filename}"
        }
        
    # Check if batch file success
    if os.path.exists(os.path.join(job_dir, "success.json")):
        with open(os.path.join(job_dir, "success.json"), "r", encoding="utf-8") as f:
            results_data = json.load(f)
        return {
            "success": True,
            "status": "completed",
            "results": results_data,
            "download_url": f"/static/conversions/{job_id}_archive.zip"
        }
        
    return {"success": True, "status": "processing"}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("app:app", host="0.0.0.0", port=8000, reload=True)
