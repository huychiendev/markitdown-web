"""
Subprocess worker for file conversion.
Runs in an isolated process so OOM kills only this process, not the main server.
Usage: python worker.py <job_id> <job_dir> <file_path> <filename>
"""
import sys
import os
import re
import json
import urllib.parse


def _escape_cell(text):
    return text.replace('|', '\\|').replace('\n', ' ')


def _extract_vba_macros(file_path):
    try:
        from oletools.olevba import VBA_Parser
        vba_parser = VBA_Parser(file_path)
        if not vba_parser.detect_vba_macros():
            vba_parser.close()
            return ""
        parts = ["\n## VBA Macros\n"]
        for (filename, stream_path, vba_filename, vba_code_chunk) in vba_parser.extract_macros():
            parts.append(f"### Module: {vba_filename}\n```vba\n{vba_code_chunk}\n```\n")
        vba_parser.close()
        return "\n".join(parts)
    except Exception as e:
        print(f"Error extracting VBA: {e}")
        return ""


def convert_excel_with_formulas(job_id, job_dir, file_path):
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter
    wb_data = load_workbook(file_path, data_only=True, read_only=False)
    wb_formula = load_workbook(file_path, data_only=False, read_only=False)
    parts = []

    for name in wb_formula.sheetnames:
        ws_d, ws_f = wb_data[name], wb_formula[name]
        
        images_by_cell = {}
        img_counter = 0
        max_img_row = -1
        max_img_col = -1
        
        for image in getattr(ws_f, '_images', []):
            try:
                if not hasattr(image, 'anchor') or not hasattr(image.anchor, '_from'):
                    continue
                
                row_idx = image.anchor._from.row
                col_idx = image.anchor._from.col
                
                if row_idx > max_img_row: max_img_row = row_idx
                if col_idx > max_img_col: max_img_col = col_idx
                
                img_counter += 1
                ext = getattr(image, 'format', 'png') or 'png'
                save_name = f"{name}_img_{img_counter}.{ext}"
                save_name = re.sub(r"[^\w\.-]", "_", save_name)
                img_path = os.path.join(job_dir, save_name)
                
                if callable(getattr(image, '_data', None)):
                    blob = image._data()
                else:
                    blob = image._data
                
                with open(img_path, "wb") as f:
                    f.write(blob)
                    
                encoded_filename = urllib.parse.quote(save_name)
                img_url = f"/static/conversions/{job_id}/{encoded_filename}"
                md_img = f"![{save_name}]({img_url})"
                
                if (row_idx, col_idx) not in images_by_cell:
                    images_by_cell[(row_idx, col_idx)] = []
                images_by_cell[(row_idx, col_idx)].append(md_img)
            except Exception as e:
                print(f"Error extracting image from Excel: {e}")
                
        actual_max_row = max(ws_d.max_row, max_img_row + 1)
        actual_max_col = max(ws_d.max_column, max_img_col + 1)

        sheet_title = name
        if getattr(ws_f, 'sheet_state', 'visible') != 'visible':
            sheet_title += " (sheet ẩn)"
        parts.append(f"## {sheet_title}\n")
        
        rows = []
        for r_idx, (row_d, row_f) in enumerate(zip(
            ws_d.iter_rows(max_row=actual_max_row, max_col=actual_max_col, values_only=True),
            ws_f.iter_rows(max_row=actual_max_row, max_col=actual_max_col, values_only=False)
        )):
            row_data = []
            for c_idx, (val, cell_f) in enumerate(zip(row_d, row_f)):
                raw = cell_f.value if cell_f else None
                if isinstance(raw, str) and raw.startswith('='):
                    text = f"{val if val is not None else ''} (`{raw}`)"
                else:
                    text = str(val) if val is not None else ''
                    
                imgs = images_by_cell.get((r_idx, c_idx), [])
                if imgs:
                    text = text + " " + " ".join(imgs) if text else " ".join(imgs)
                    
                row_data.append(_escape_cell(text.strip()))
                
            if any(c.strip() != '' for c in row_data):
                rows.append(row_data)

        if not rows:
            continue
        ncols = max(len(r) for r in rows)
        if ncols == 0:
            continue
            
        header_row = [get_column_letter(i + 1) for i in range(ncols)]
        parts.append('| ' + ' | '.join(header_row) + ' |')
        parts.append('| ' + ' | '.join(['---'] * ncols) + ' |')
        for row in rows:
            parts.append('| ' + ' | '.join((row + [''] * ncols)[:ncols]) + ' |')
        parts.append('')

    wb_data.close()
    wb_formula.close()
    return '\n'.join(parts)


def convert_file(job_id, job_dir, file_path, filename):
    """Core conversion logic -- runs in isolated process."""
    markdown_text = ""

    if filename.lower().endswith(('.xlsx', '.xlsm', '.xltx', '.xltm')):
        markdown_text = convert_excel_with_formulas(job_id, job_dir, file_path)
    else:
        # Only import MarkItDown here -- keeps main server process lean
        from markitdown import MarkItDown
        md = MarkItDown()
        result = md.convert(file_path)
        markdown_text = result.text_content
        markdown_text = re.sub(r'\bNaN\b', '', markdown_text)
        markdown_text = re.sub(r'Unnamed:\s*\d+', '', markdown_text)

    # VBA macros
    office_exts = ('doc', 'docx', 'docm', 'dot', 'dotx', 'dotm',
                   'xls', 'xlsx', 'xlsm', 'xlsb', 'xlt', 'xltx', 'xltm',
                   'ppt', 'pptx', 'pptm', 'pot', 'potx', 'potm', 'pps', 'ppsx', 'ppsm')
    if filename.lower().endswith(office_exts):
        vba_text = _extract_vba_macros(file_path)
        if vba_text:
            markdown_text += vba_text

    # PPTX images
    if filename.lower().endswith('.pptx'):
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            img_counter = {}
            for slide in prs.slides:
                sorted_shapes = sorted(
                    slide.shapes,
                    key=lambda s: (
                        float("-inf") if not s.top else s.top,
                        float("-inf") if not s.left else s.left,
                    ),
                )
                for shape in sorted_shapes:
                    if not (shape.shape_type == 13 or (shape.shape_type == 14 and hasattr(shape, "image"))):
                        continue
                    try:
                        blob = shape.image.blob
                        ext = shape.image.ext
                    except Exception:
                        continue

                    placeholder_name = re.sub(r"\W", "", shape.name)
                    placeholder = f"{placeholder_name}.jpg"
                    count = img_counter.get(placeholder_name, 0)
                    img_counter[placeholder_name] = count + 1
                    save_name = f"{placeholder_name}_{count}.{ext}" if count > 0 else f"{placeholder_name}.{ext}"

                    img_path = os.path.join(job_dir, save_name)
                    with open(img_path, "wb") as img_f:
                        img_f.write(blob)

                    encoded_filename = urllib.parse.quote(save_name)
                    img_url = f"/static/conversions/{job_id}/{encoded_filename}"
                    markdown_text = markdown_text.replace(f"]({placeholder})", f"]({img_url})", 1)
        except Exception as e:
            print(f"Error extracting PPTX images: {e}")

    return markdown_text


def main():
    if len(sys.argv) < 5:
        print("Usage: python worker.py <job_id> <job_dir> <file_path> <filename> [--batch]")
        sys.exit(1)

    job_id, job_dir, file_path, filename = sys.argv[1], sys.argv[2], sys.argv[3], sys.argv[4]
    is_batch = len(sys.argv) > 5 and sys.argv[5] == "--batch"

    try:
        markdown_text = convert_file(job_id, job_dir, file_path, filename)

        md_filename = f"{os.path.splitext(filename)[0]}.md"
        md_filepath = os.path.join(job_dir, md_filename)
        os.makedirs(os.path.dirname(md_filepath), exist_ok=True)
        
        with open(md_filepath, "w", encoding="utf-8") as f:
            f.write(markdown_text)

        if not is_batch:
            with open(os.path.join(job_dir, "success.txt"), "w") as f:
                f.write(md_filename)

    except Exception as e:
        if not is_batch:
            with open(os.path.join(job_dir, "error.txt"), "w", encoding="utf-8") as f:
                f.write(str(e))
        else:
            print(f"Error processing {filename}: {e}", file=sys.stderr)
            sys.exit(1)
    finally:
        if os.path.exists(file_path):
            os.remove(file_path)


if __name__ == "__main__":
    main()
